import win32com.client
import os
import openai
from dotenv import load_dotenv
import time
import logging
import re
from pptx import Presentation # python-pptx 임포트
from pptx.util import Inches, Pt

# --- Configuration & Prerequisite Functions ---
# (load_dotenv, openai.api_key setup, logging setup)
# (get_user_input 함수는 그대로 사용)
# (VBA_ENTRY_POINT_MACRO 상수 정의)

CUR_PATH = os.path.dirname(__file__)

load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")
if not openai.api_key:
    raise ValueError("OpenAI API key not found. Please set it in the .env file.")

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

VBA_ENTRY_POINT_MACRO = "CreatePresentationFromData"

def get_user_input() -> dict:
    """사용자 입력을 받습니다."""
    logging.info("Collecting user input...")
    content = {}
    content['title'] = input("Enter the main title for the presentation: ")
    content['slides'] = []
    while True:
        slide_title = input("Enter title for a new slide (or leave blank to finish): ")
        if not slide_title:
            break
        slide_points = []
        print("Enter bullet points for this slide (leave blank to finish):")
        while True:
            point = input("- ")
            if not point:
                break
            slide_points.append(point)
        content['slides'].append({'title': slide_title, 'points': slide_points})
    logging.info(f"User input collected: {content}")
    return content


# --- Step 1: Analyze PPTX ---
def analyze_pptx_structure(pptx_path: str) -> dict | None:
    """
    Uses python-pptx to analyze the structure (layouts) of a .pptx file.
    """
    if not os.path.exists(pptx_path):
        logging.error(f"PPTX template file not found: {pptx_path}")
        return None
    if not pptx_path.lower().endswith(".pptx"):
        logging.warning(f"File is not a .pptx file: {pptx_path}")
        return None # Only analyze .pptx

    logging.info(f"Analyzing PPTX structure: {pptx_path}")
    try:
        prs = Presentation(pptx_path)
        layout_info = {}
        # Iterate through slide layouts in the master
        for i, layout in enumerate(prs.slide_layouts):
            # Try to get placeholder info (this is basic, might need refinement)
            placeholders = []
            try:
                for shape in layout.placeholders:
                     # Get placeholder type name if possible, default to shape name
                     ph_type = shape.name
                     if shape.is_placeholder:
                         ph_type = shape.placeholder_format.type # Get type enum
                     placeholders.append(str(ph_type)) # Store type as string
            except Exception: # Handle layouts with potentially complex placeholders
                pass # Keep placeholders list potentially empty for this layout
            layout_info[layout.name] = {'index': i, 'placeholders': placeholders}

        logging.info(f"Found layouts: {list(layout_info.keys())}")
        # Return analysis results including the path for theme application
        return {'layouts': layout_info, 'theme_file': os.path.abspath(pptx_path)}
    except Exception as e:
        logging.error(f"Failed to analyze pptx structure using python-pptx: {e}", exc_info=True)
        return None

# --- Step 3 & 4: Generate VBA with AI (Modified Prompt) ---
def generate_vba_based_on_analysis(analyzed_data: dict, user_data: dict) -> str | None:
    """
    Generates VBA code using AI, informed by the PPTX analysis and user data.
    Theme application instruction is removed from VBA, handled by pywin32 instead.
    """
    logging.info("Generating VBA based on PPTX analysis and user data...")
    layout_summary = "\n".join([f"- '{name}' (Layout Index {info['index']})" for name, info in analyzed_data.get('layouts', {}).items()])
    theme_file = analyzed_data.get('theme_file', 'the_original.pptx') # For context only

    prompt = f"""
    Generate VBA code for Microsoft PowerPoint to create a presentation based on user data.

    **Context from PPTX Analysis:**
    The generated presentation will use a design based on '{os.path.basename(theme_file)}'.
    The original presentation includes the following slide layouts (use these as a guide for choosing layout indices):
    {layout_summary if layout_summary else "Standard layouts presumed."}
    NOTE: The theme/design will be applied *before* this VBA code runs. Do NOT include ApplyTheme or design changes in the generated VBA code.

    **User Data to Incorporate:**
    ```json
    {user_data}
    ```

    **Instructions for VBA Generation:**
    1. Create a main public subroutine named `{VBA_ENTRY_POINT_MACRO}`. This sub should operate on the `ActivePresentation`.
    2. Use the provided JSON user data to add slides.
    3. When adding slides (e.g., `ActivePresentation.Slides.AddSlide(Index, CustomLayout)`):
        - Determine the appropriate `CustomLayout` object using the layout index. Use `ActivePresentation.SlideMaster.CustomLayouts(LayoutIndex)`.
        - Choose a `LayoutIndex` based on the 'Context from PPTX Analysis' above. For example, if adding a slide with a title and bullet points, use the index corresponding to a 'Title and Content' layout (often index 1, but refer to the list). If unsure, use a common default like 1.
        - The `Index` argument for `AddSlide` determines the position of the new slide (e.g., `ActivePresentation.Slides.Count + 1` to add at the end).
    4. Populate slide titles and bullet points from the JSON data onto the placeholders of the newly added slides. Use standard placeholder access methods (e.g., `newSlide.Shapes.Title`, `newSlide.Shapes.Placeholders(2)` for body). Handle potential errors if placeholders don't exist as expected (e.g., using `On Error Resume Next`).
    5. Ensure the VBA code is self-contained, uses `Option Explicit`, and is syntactically correct. Avoid adding comments unless necessary for complex logic.
    6. Output ONLY the raw VBA code, without any introduction, explanation, or markdown formatting.

    **Generated VBA Code:**
    """
    try:
        # Using the chat completions endpoint (recommended)
        response = openai.chat.completions.create(
            model="gpt-4-turbo",  # Or another suitable model
            messages=[
                {"role": "system", "content": "You are an expert assistant that generates Microsoft PowerPoint VBA code according to specific instructions."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.4, # Lower temperature for more predictable code
            max_tokens=2500  # Adjust as needed
        )

        generated_code = response.choices[0].message.content.strip()
        # Clean potential markdown fences
        generated_code = re.sub(r'^```vb\s*|\s*```$', '', generated_code, flags=re.MULTILINE).strip()

        logging.info("AI generated VBA code successfully based on analysis.")
        logging.debug(f"Generated VBA (start):\n{generated_code[:300]}")
        return generated_code

    except Exception as e:
        logging.error(f"Error calling OpenAI API: {e}", exc_info=True)
        return None


# --- Step 5: Create Presentation (Modified to Apply Theme) ---
def create_ppt_with_vba_and_theme(generated_vba: str, output_path: str, entry_point_macro: str, theme_pptx_path: str | None):
    """
    Creates a new PPT, applies theme from the source PPTX,
    injects the generated VBA, runs it, and saves the presentation.
    """
    powerpoint = None
    presentation = None
    success = False

    # Ensure output path has .pptm extension
    if not output_path.lower().endswith(".pptm"):
        output_path = os.path.splitext(output_path)[0] + ".pptm"
        logging.info(f"Output filename adjusted to: {output_path}")

    try:
        logging.info("Creating new PowerPoint instance for generation...")
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        powerpoint.Visible = True # Make visible for debugging
        presentation = powerpoint.Presentations.Add()
        logging.info("New presentation created.")
        abs_output_path = os.path.abspath(output_path)

        # --- Apply Theme ---
        if theme_pptx_path and os.path.exists(theme_pptx_path):
            abs_theme_path = os.path.abspath(theme_pptx_path)
            logging.info(f"Attempting to apply theme from: {abs_theme_path}")
            try:
                presentation.ApplyTheme(abs_theme_path)
                # Allow a moment for theme application
                time.sleep(1)
                logging.info("Theme applied successfully.")
            except Exception as e_theme:
                logging.error(f"Failed to apply theme: {e_theme}. Continuing without custom theme.")
        elif theme_pptx_path:
             logging.warning(f"Theme file not found, skipping ApplyTheme: {theme_pptx_path}")
        else:
            logging.info("No theme file provided, using default theme.")


        # --- Inject and Run VBA ---
        logging.info("Adding VBA module...")
        # Add module after potential theme application
        vb_project = presentation.VBProject
        vb_module = vb_project.VBComponents.Add(1) # 1 = vbext_ct_StdModule
        module_name = "GeneratedAIModule"
        vb_module.Name = module_name
        logging.info(f"Module '{module_name}' added.")

        logging.info("Injecting generated VBA code...")
        vb_module.CodeModule.AddFromString(generated_vba)
        logging.info("VBA code injected.")

        full_macro_path = f"{module_name}.{entry_point_macro}"
        logging.info(f"Attempting to run macro: {full_macro_path}")
        time.sleep(2) # Give PowerPoint a moment

        try:
             powerpoint.Run(full_macro_path)
             logging.info(f"Successfully executed macro: {full_macro_path}")
             success = True
        except Exception as e_run:
            logging.error(f"ERROR RUNNING MACRO '{full_macro_path}': {e_run}")
            logging.error("Potential issues: Macro name, VBA syntax error, security settings, PP instability.")
            success = False # Mark as failed

        # --- Save ---
        if success: # Only save if macro execution seemed successful (or adjust logic)
            logging.info(f"Saving presentation to: {abs_output_path}")
            # FileFormat 25 = ppSaveAsMacroEnabledPresentation
            presentation.SaveAs(abs_output_path, FileFormat=25)
            logging.info("Presentation saved successfully.")
        else:
            logging.warning("Macro execution failed or was skipped, presentation not saved automatically.")
            # Optionally save anyway for debugging:
            # try:
            #     debug_path = os.path.splitext(abs_output_path)[0] + "_debug.pptm"
            #     presentation.SaveAs(debug_path, FileFormat=25)
            #     logging.info(f"Saved debug presentation (with potentially non-working VBA) to: {debug_path}")
            # except Exception as e_save_debug:
            #     logging.error(f"Could not save debug presentation: {e_save_debug}")


    except Exception as e:
        logging.error(f"Error during PowerPoint creation/automation: {e}", exc_info=True)
        success = False
    finally:
        # Close presentation (if open)
        if presentation:
            try:
                # Close without saving changes again
                presentation.Close()
            except Exception: pass
        # Quit PowerPoint application
        if powerpoint:
            try:
                powerpoint.Quit()
                time.sleep(1)
            except Exception: pass
        # Clean up COM objects (optional, Python's GC usually handles it)
        presentation = None
        powerpoint = None

    return success

# --- Main Execution Logic ---
if __name__ == "__main__":
    template_file = os.path.join(CUR_PATH, "sample_ppt.pptx")
    output_file = os.path.join(CUR_PATH, "sample_ppt.pptm")

    # Basic path validation
    template_file = os.path.abspath(template_file)
    output_file = os.path.abspath(output_file) # Will be adjusted to .pptm later

    if not template_file.lower().endswith(".pptx"):
        logging.error("This workflow currently requires a .pptx file as the template input.")
        exit()
    if not os.path.exists(template_file):
         logging.error(f"Template file not found: {template_file}")
         exit()

    # --- Workflow Steps ---
    # 1. Analyze PPTX structure
    logging.info("--- Step 1: Analyzing PPTX Structure ---")
    analyzed_info = analyze_pptx_structure(template_file)
    if not analyzed_info:
        logging.error("Failed to analyze the source .pptx file.")
        exit()

    # 2. Get User Input
    logging.info("--- Step 2: Getting User Input ---")
    user_content = get_user_input()

    # 3. Generate VBA based on Analysis and User Input
    logging.info("--- Step 3: Generating VBA with AI ---")
    generated_vba_code = generate_vba_based_on_analysis(analyzed_info, user_content)
    if not generated_vba_code:
        logging.error("Failed to generate VBA code using AI.")
        exit()

    # 4. Create PPT (Apply Theme, Inject & Run VBA)
    logging.info("--- Step 4: Creating PowerPoint Presentation ---")
    creation_success = create_ppt_with_vba_and_theme(
        generated_vba_code,
        output_file, # Path will be adjusted to .pptm inside the function
        VBA_ENTRY_POINT_MACRO,
        analyzed_info.get('theme_file') # Pass the original pptx path for theme application
    )

    if creation_success:
        final_output_path = os.path.splitext(output_file)[0] + ".pptm"
        logging.info(f"--- Process Complete: Presentation saved to {final_output_path} ---")
    else:
        logging.error("--- Process Failed: Presentation creation failed. Check logs above. ---")

    logging.info("Script finished.")